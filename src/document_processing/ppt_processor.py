class PptProcessor:
    def __init__(self, file_path):
        self.file_path = file_path

    def convert_to_images(self):
        from pptx import Presentation
        from PIL import Image
        import io

        presentation = Presentation(self.file_path)
        images = []

        for slide in presentation.slides:
            img_stream = io.BytesIO()
            slide.shapes._spTree.write(img_stream)
            img_stream.seek(0)
            img = Image.open(img_stream)
            images.append(img)

        return images

    def save_images(self, images, output_folder):
        import os

        if not os.path.exists(output_folder):
            os.makedirs(output_folder)

        for i, img in enumerate(images):
            img.save(os.path.join(output_folder, f'slide_{i + 1}.png'), 'PNG')